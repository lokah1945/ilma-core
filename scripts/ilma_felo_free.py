#!/usr/bin/env python3
"""
ILMA FELO FREE — 100% Gratis Felo AI Replacements
==================================================
Replace all Felo paid API calls with native free alternatives.

FREE REPLACEMENTS MAP:
=======================
Felo Paid (REPLACE)        → ILMA Free Native
──────────────────────────────────────────────────
felo-search                 → web_search tool + browser_navigate
felo-x-search              → browser_navigate + BeautifulSoup
felo-slides                → python-pptx code generation
felo-mindmap               → HTML/Mermaid mindmap
felo-web-fetch             → browser_navigate + extract text
felo-youtube-subtitling    → yt-dlp OR browser fetch
felo-superagent            → MiniMax-M2.7 with session memory
felo-livedoc               → file-based knowledge base
felo-content-to-slides     → browser_navigate + python-pptx
felo-twitter-writer        → browser scrape + MiniMax analysis

RULES:
- NEVER call Felo API (openapi.felo.ai)
- NEVER use FELO_API_KEY for these replacements
- Always log: "Using FREE native method"
- Fallback: try Felo API only if native fails AND key is set

AUTHOR: ILMA Hermes Agent
DATE: 2026-05-08
"""

import subprocess
import json
import re
import os
import sys
from datetime import datetime

# =====================================================================
# CONFIG
# =====================================================================

SKILL_LOG = "/root/.hermes/profiles/ilma/docs/ILMA_felo_free_log.md"
OUTPUT_DIR = "/root/.hermes/profiles/ilma/output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# =====================================================================
# LOGGING
# =====================================================================

def log(skill: str, status: str, detail: str = ""):
    ts = datetime.now().strftime("%Y-%m-%d %H:%M")
    entry = f"| {ts} | {skill} | {status} | {detail} |\n"
    try:
        with open(SKILL_LOG, "a") as f:
            f.write(entry)
    except IOError:
        pass
    print(f"[{skill}] {status}: {detail}")

# =====================================================================
# 1. NATIVE SEARCH (Replace: felo-search)
# =====================================================================

def native_search(query: str, limit: int = 10) -> dict:
    """
    Replace felo-search API with native web search.
    Uses: requests + BeautifulSoup (no API key needed)
    
    Free sources:
    - DuckDuckGo (html mode)
    - Bing (via serper.dev free tier)
    - Wikipedia API
    
    Returns: dict with 'results', 'answer', 'sources'
    """
    import requests
    from bs4 import BeautifulSoup
    
    log("native_search", "ACTIVE", f"Query: {query[:50]}")
    
    results = []
    sources = []
    
    # Source 1: DuckDuckGo HTML
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36",
            "Accept-Language": "en-US,en;q=0.9"
        }
        url = f"https://html.duckduckgo.com/html/?q={requests.utils.quote(query)}"
        resp = requests.get(url, headers=headers, timeout=10)
        soup = BeautifulSoup(resp.text, "html.parser")
        
        for result in soup.select(".result")[:limit]:
            title_el = result.select_one(".result__title a")
            snippet_el = result.select_one(".result__snippet")
            if title_el:
                title = title_el.get_text(strip=True)
                link = title_el.get("href", "")
                snippet = snippet_el.get_text(strip=True) if snippet_el else ""
                results.append({"title": title, "url": link, "snippet": snippet})
                sources.append(link)
    except Exception as e:
        log("native_search", "WARN", f"DuckDuckGo failed: {e}")
    
    # Source 2: Wikipedia API
    try:
        wiki_url = f"https://en.wikipedia.org/w/api.php"
        params = {
            "action": "opensearch",
            "search": query,
            "limit": 5,
            "format": "json"
        }
        resp = requests.get(wiki_url, params=params, timeout=10)
        wiki_data = resp.json()
        
        if len(wiki_data) >= 4:
            for i, title in enumerate(wiki_data[1][:5]):
                if i < len(wiki_data[2]):
                    snippet = wiki_data[2][i]
                else:
                    snippet = ""
                url = wiki_data[3][i] if i < len(wiki_data[3]) else ""
                if title and url:
                    results.append({
                        "title": title,
                        "url": url,
                        "snippet": snippet,
                        "source": "Wikipedia"
                    })
                    sources.append(url)
    except Exception as e:
        log("native_search", "WARN", f"Wikipedia failed: {e}")
    
    answer = f"Search results for '{query}' found {len(results)} items."
    
    log("native_search", "DONE", f"Found {len(results)} results")
    
    return {
        "status": "ok",
        "query": query,
        "results": results[:limit],
        "sources": sources[:limit],
        "answer": answer,
        "method": "FREE_NATIVE (DuckDuckGo + Wikipedia)"
    }


# =====================================================================
# 2. NATIVE X/TWITTER SEARCH (Replace: felo-x-search)
# =====================================================================

def native_x_search(query: str, limit: int = 20) -> dict:
    """
    Replace felo-x-search with native browser scraping.
    
    Free methods (in order of preference):
    1. Nitter instances (RSS/API, no auth) — deprecated Twitter disable
    2. Brave Search API (free tier)
    3. Browser scrape via playwright
    4. Alternative: snscrape (if installed)
    
    Returns: dict with 'tweets', 'users'
    """
    from bs4 import BeautifulSoup
    import requests
    
    log("native_x_search", "ACTIVE", f"Query: {query[:50]}")
    
    tweets = []
    
    # Method 1: Nitter (free Twitter archive)
    nitter_instances = [
        "nitter.poast.org",
        "nitter.privacydev.net", 
        "nitter.fdn.fr",
        "nitter.unixfox.eu",
    ]
    
    for instance in nitter_instances:
        try:
            url = f"https://{instance}/search?f=tweets&q={requests.utils.quote(query)}"
            headers = {"User-Agent": "Mozilla/5.0"}
            resp = requests.get(url, headers=headers, timeout=10)
            if resp.status_code != 200:
                continue
            
            soup = BeautifulSoup(resp.text, "html.parser")
            for tweet in soup.select(".timeline-item")[:limit]:
                content = tweet.get_text(strip=True)
                if content and len(content) > 20:
                    tweets.append({
                        "content": content,
                        "source": f"Nitter/{instance}",
                        "url": url
                    })
            if tweets:
                log("native_x_search", "DONE", f"Nitter worked on {instance}, {len(tweets)} tweets")
                break
        except Exception as e:
            continue
    
    # Method 2: Brave Search API (if key available)
    if not tweets:
        try:
            import os
            brave_key = os.environ.get("BRAVE_SEARCH_API_KEY", "")
            if brave_key:
                url = "https://api.search.brave.com/res/v1/search"
                headers = {
                    "X-Subscription-Token": brave_key,
                    "Accept": "application/json"
                }
                params = {"q": query, "count": limit}
                resp = requests.get(url, headers=headers, params=params, timeout=10)
                data = resp.json()
                
                if "web" in data and "results" in data["web"]:
                    for item in data["web"]["results"][:limit]:
                        tweets.append({
                            "content": item.get("title", ""),
                            "snippet": item.get("description", ""),
                            "url": item.get("url", ""),
                            "source": "Brave Search"
                        })
        except Exception:
            pass
    
    if not tweets:
        log("native_x_search", "WARN", "No results from any method")
    
    return {
        "status": "ok" if tweets else "partial",
        "query": query,
        "tweets": tweets,
        "count": len(tweets),
        "method": "FREE_NATIVE (Nitter/Brave)",
        "note": "X/Twitter API officially disabled free access 2024. Native methods have limited coverage."
    }


# =====================================================================
# 3. NATIVE PPT/SLIDES (Replace: felo-slides)
# =====================================================================

def native_slides(topic: str, slides: list = None, output_path: str = None) -> dict:
    """
    Replace felo-slides API with python-pptx code generation.
    
    Usage:
        native_slides("AI Trends", [
            {"title": "Slide 1", "content": "Bullet points..."},
            {"title": "Slide 2", "content": "More content..."},
        ])
    """
    log("native_slides", "ACTIVE", f"Topic: {topic}")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
    except ImportError:
        log("native_slides", "ERROR", "python-pptx not installed. Installing...")
        result = subprocess.run(
            [sys.executable, "-m", "pip", "install", "python-pptx", "-q"],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode != 0:
            log("native_slides", "ERROR", "pip install failed")
            return {"status": "error", "message": "Could not install python-pptx"}
        from pptx import Presentation
    
    if output_path is None:
        output_path = os.path.join(OUTPUT_DIR, f"slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    if slides is None:
        slides = [
            {"title": topic, "content": "Introduction"},
            {"title": "Overview", "content": "Key points"},
            {"title": "Details", "content": "Deep dive"},
            {"title": "Conclusion", "content": "Summary"},
        ]
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = topic
    p.font.size = Pt(44)
    p.font.bold = True
    
    # Content slides
    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", f"Slide {i+2}")
        p.font.size = Pt(32)
        p.font.bold = True
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("content", "")
        p.font.size = Pt(18)
    
    prs.save(output_path)
    log("native_slides", "DONE", f"Saved: {output_path}")
    
    return {
        "status": "ok",
        "topic": topic,
        "slide_count": len(slides) + 1,
        "output_path": output_path,
        "method": "FREE_NATIVE (python-pptx)"
    }


# =====================================================================
# 4. NATIVE MINDMAP (Replace: felo-mindmap)
# =====================================================================

def native_mindmap(topic: str, branches: list = None, output_path: str = None) -> dict:
    """
    Replace felo-mindmap with HTML/Mermaid.js mindmap.
    
    Output: Interactive HTML mindmap
    """
    log("native_mindmap", "ACTIVE", f"Topic: {topic}")
    
    if output_path is None:
        output_path = os.path.join(OUTPUT_DIR, f"mindmap_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html")
    
    if branches is None:
        branches = ["Overview", "Key Points", "Details", "Summary"]
    
    # Build Mermaid markdown
    mermaid_lines = [f"mindmap", f"  root({topic})"]
    for i, branch in enumerate(branches):
        mermaid_lines.append(f"    {i+1}[{branch}]")
        mermaid_lines.append(f"      {i+1}1[Sub point 1]")
        mermaid_lines.append(f"      {i+1}2[Sub point 2]")
    
    mermaid_md = "\n".join(mermaid_lines)
    
    html = f"""<!DOCTYPE html>
<html>
<head>
    <title>Mindmap: {topic}</title>
    <script src="https://cdn.jsdelivr.net/npm/mermaid/dist/mermaid.min.js"></script>
    <style>
        body {{ font-family: Arial, sans-serif; background: #1a1a2e; color: #eee; padding: 20px; }}
        .mermaid {{ background: #16213e; border-radius: 10px; padding: 20px; }}
        h1 {{ color: #e94560; }}
        .meta {{ color: #888; font-size: 14px; }}
    </style>
</head>
<body>
    <h1>🧠 Mindmap: {topic}</h1>
    <p class="meta">Generated by ILMA FELO FREE | {datetime.now().strftime('%Y-%m-%d %H:%M')}</p>
    <div class="mermaid">
{mermaid_md}
    </div>
    <script>mermaid.initialize({{ theme: 'dark', mindmap: {{ padding: 20 }} }});</script>
</body>
</html>"""
    
    with open(output_path, "w") as f:
        f.write(html)
    
    log("native_mindmap", "DONE", f"Saved: {output_path}")
    
    return {
        "status": "ok",
        "topic": topic,
        "output_path": output_path,
        "mermaid_code": mermaid_md,
        "method": "FREE_NATIVE (Mermaid.js HTML)"
    }


# =====================================================================
# 5. NATIVE WEB FETCH (Replace: felo-web-fetch)
# =====================================================================

def native_web_fetch(url: str, format: str = "markdown") -> dict:
    """
    Replace felo-web-fetch API with native browser + BeautifulSoup.
    
    format: 'markdown', 'html', 'text'
    """
    from bs4 import BeautifulSoup
    import requests
    
    log("native_web_fetch", "ACTIVE", f"URL: {url[:60]}")
    
    headers = {
        "User-Agent": "Mozilla/5.0 (compatible; ILMA/1.0)"
    }
    
    try:
        resp = requests.get(url, headers=headers, timeout=15, allow_redirects=True)
        resp.raise_for_status()
    except Exception as e:
        log("native_web_fetch", "ERROR", str(e))
        return {"status": "error", "message": str(e)}
    
    if format == "html":
        content = resp.text
    else:
        soup = BeautifulSoup(resp.text, "html.parser")
        
        # Remove script/style/nav/footer
        for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
            tag.decompose()
        
        if format == "text":
            content = soup.get_text(separator="\n", strip=True)
        else:  # markdown approximation
            content = soup.get_text(separator="\n\n", strip=True)
    
    # Extract title
    title_match = re.search(r'<title>([^<]+)</title>', resp.text)
    title = title_match.group(1) if title_match else url
    
    log("native_web_fetch", "DONE", f"Extracted {len(content)} chars from {title}")
    
    return {
        "status": "ok",
        "url": url,
        "title": title,
        "content": content[:50000],  # limit
        "format": format,
        "method": "FREE_NATIVE (requests + BeautifulSoup)"
    }


# =====================================================================
# 6. NATIVE YOUTUBE SUBTITLES (Replace: felo-youtube-subtitling)
# =====================================================================

def native_youtube_subtitles(video_id: str = None, video_url: str = None) -> dict:
    """
    Replace felo-youtube-subtitling with native extraction.
    
    Methods:
    1. yt-dlp (if installed) — BEST
    2. YouTube transcript API (no auth) — via youtube-transcript-api
    3. Browser fetch — fallback
    
    video_id: YouTube video ID (e.g., "dQw4w9WgXcQ")
    video_url: Full YouTube URL
    """
    import requests
    
    log("native_youtube_subtitles", "ACTIVE", f"Video: {video_id or video_url}")
    
    if video_url:
        # Extract video ID
        match = re.search(r'(?:v=|/v/|youtu\.be/)([a-zA-Z0-9_-]{11})', video_url)
        if match:
            video_id = match.group(1)
    
    if not video_id:
        return {"status": "error", "message": "No video ID found"}
    
    subtitles = []
    
    # Method 1: youtube-transcript-api (no auth needed)
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        transcript = transcript_list.find_transcript(['en'])
        data = transcript.fetch()
        
        subtitle_text = "\n".join([item['text'] for item in data])
        subtitles.append({
            "language": "en",
            "is_generated": transcript.is_generated,
            "text": subtitle_text
        })
        log("native_youtube_subtitles", "DONE", "youtube-transcript-api succeeded")
    except ImportError:
        log("native_youtube_subtitles", "WARN", "youtube-transcript-api not installed")
        try:
            result = subprocess.run(
                [sys.executable, "-m", "pip", "install", "youtube-transcript-api", "-q"],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0:
                from youtube_transcript_api import YouTubeTranscriptApi
                transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
                transcript = transcript_list.find_transcript(['en'])
                data = transcript.fetch()
                subtitle_text = "\n".join([item['text'] for item in data])
                subtitles.append({
                    "language": "en",
                    "is_generated": transcript.is_generated,
                    "text": subtitle_text
                })
        except Exception as e:
            log("native_youtube_subtitles", "WARN", f"Method 1 failed: {e}")
    except Exception as e:
        log("native_youtube_subtitles", "WARN", f"Method 1 failed: {e}")
    
    # Method 2: yt-dlp (if installed)
    if not subtitles:
        try:
            result = subprocess.run(
                ["yt-dlp", "--list-subs", video_id or video_url],
                capture_output=True, text=True, timeout=10
            )
            if result.returncode == 0 and "--en" in result.stdout:
                # Has English subs
                result = subprocess.run(
                    ["yt-dlp", "--write-sub", "--sub-lang", "en", 
                     "--skip-download", "--output", "/tmp/yt_sub", 
                     video_id or video_url],
                    capture_output=True, text=True, timeout=30
                )
                sub_file = "/tmp/yt_sub.en.vtt"
                if os.path.exists(sub_file):
                    with open(sub_file) as f:
                        subtitles.append({"language": "en", "text": f.read()})
        except IOError:
            pass
    
    return {
        "status": "ok" if subtitles else "partial",
        "video_id": video_id,
        "subtitles": subtitles,
        "method": "FREE_NATIVE (youtube-transcript-api / yt-dlp)"
    }


# =====================================================================
# 7. NATIVE SUPERAGENT (Replace: felo-superagent)
# =====================================================================

def native_superagent(query: str, context: str = "", session_id: str = None) -> dict:
    """
    Replace felo-superagent with MiniMax-M2.7 (current model).
    
    Session management via file-based memory.
    Context: previous turns stored in session file.
    
    Returns: response text + session update
    """
    import requests
    
    log("native_superagent", "ACTIVE", f"Query: {query[:50]}")
    
    # Load session context
    if session_id is None:
        session_id = datetime.now().strftime("%Y%m%d")
    
    session_file = os.path.join(OUTPUT_DIR, f"session_{session_id}.json")
    
    try:
        with open(session_file) as f:
            session = json.load(f)
    except ValueError:
        session = {"turns": [], "query_count": 0}
    
    # Build conversation history
    messages = []
    for turn in session.get("turns", [])[-10:]:  # last 10 turns
        messages.append({"role": "user", "content": turn["user"]})
        messages.append({"role": "assistant", "content": turn["assistant"]})
    
    if context:
        messages.append({"role": "system", "content": context})
    
    messages.append({"role": "user", "content": query})
    
    # Call MiniMax-M2.7 (current model)
    try:
        minimax_key = os.environ.get("MINIMAX_API_KEY", "")
        if not minimax_key:
            # Try to get from credential file
            cred_file = "/root/credential/api_key.json"
            if os.path.exists(cred_file):
                with open(cred_file) as f:
                    creds = json.load(f)
                    if "minimax" in creds and creds["minimax"].get("keys"):
                        minimax_key = creds["minimax"]["keys"][0]
        
        if minimax_key:
            headers = {"Authorization": f"Bearer {minimax_key}"}
            payload = {
                "model": "MiniMax-M2.7",
                "messages": messages,
                "stream": False
            }
            resp = requests.post(
                "https://api.minimax.io/v1/text/chatcompletion_v2",
                headers=headers, json=payload, timeout=60
            )
            data = resp.json()
            answer = data["choices"][0]["message"]["content"]
        else:
            # Fallback: just return query info (no actual model call)
            answer = f"[FREE MODE] Received: {query}. MiniMax API key not configured for direct calls."
            log("native_superagent", "WARN", "No MiniMax key, using free fallback")
            
    except Exception as e:
        answer = f"[FREE MODE] Error: {e}"
        log("native_superagent", "ERROR", str(e))
    
    # Save session
    session["turns"].append({"user": query, "assistant": answer})
    session["query_count"] += 1
    
    try:
        with open(session_file, "w") as f:
            json.dump(session, f, indent=2)
    except Exception:
        pass
    
    log("native_superagent", "DONE", f"Turn {session['query_count']} complete")
    
    return {
        "status": "ok",
        "query": query,
        "answer": answer,
        "session_id": session_id,
        "turn": session["query_count"],
        "method": "FREE_NATIVE (MiniMax-M2.7 session)"
    }


# =====================================================================
# 8. NATIVE KNOWLEDGE BASE (Replace: felo-livedoc)
# =====================================================================

KB_DIR = os.path.join(OUTPUT_DIR, "knowledge_base")
os.makedirs(KB_DIR, exist_ok=True)

def native_kb_create(name: str) -> dict:
    """Create a knowledge base directory."""
    kb_id = datetime.now().strftime("%Y%m%d%H%M%S")
    kb_path = os.path.join(KB_DIR, kb_id)
    os.makedirs(kb_path, exist_ok=True)
    
    meta = {"id": kb_id, "name": name, "created": datetime.now().isoformat()}
    with open(os.path.join(kb_path, "meta.json"), "w") as f:
        json.dump(meta, f)
    
    log("native_kb", "CREATED", f"KB: {name} ({kb_id})")
    return {"status": "ok", "kb_id": kb_id, "name": name, "path": kb_path}


def native_kb_search(kb_id: str, query: str, limit: int = 5) -> dict:
    """Search knowledge base (simple keyword match)."""
    import requests
    
    kb_path = os.path.join(KB_DIR, kb_id)
    if not os.path.exists(kb_path):
        return {"status": "error", "message": "KB not found"}
    
    results = []
    for fname in os.listdir(kb_path):
        if fname.startswith("."):
            continue
        fpath = os.path.join(kb_path, fname)
        if os.path.isfile(fpath):
            try:
                with open(fpath) as f:
                    content = f.read()
                if query.lower() in content.lower():
                    results.append({"file": fname, "content": content[:1000]})
            except Exception:
                pass
    
    log("native_kb", "SEARCHED", f"KB {kb_id}, query: {query[:30]}, found: {len(results)}")
    return {
        "status": "ok",
        "kb_id": kb_id,
        "query": query,
        "results": results[:limit]
    }


# =====================================================================
# 9. NATIVE CONTENT-TO-SLIDES (Replace: felo-content-to-slides)
# =====================================================================

def native_content_to_slides(url: str, topic: str = None) -> dict:
    """
    Replace felo-content-to-slides: fetch URL content + generate PPT.
    
    Steps:
    1. native_web_fetch(url) → content
    2. Extract key points (simple heuristic)
    3. native_slides(topic, slides)
    """
    log("native_content_to_slides", "ACTIVE", f"URL: {url[:60]}")
    
    # Step 1: Fetch content
    fetch_result = native_web_fetch(url, format="text")
    if fetch_result["status"] != "ok":
        return fetch_result
    
    content = fetch_result["content"]
    title = fetch_result.get("title", topic or "Content Slides")
    
    # Step 2: Extract key points (simple split by double newline)
    paragraphs = [p.strip() for p in content.split("\n\n") if p.strip() and len(p.strip()) > 30]
    
    slides = []
    for i, para in enumerate(paragraphs[:10]):  # max 10 slides
        # Take first sentence as title, rest as content
        sentences = para.split(". ")
        slide_title = sentences[0][:80] if sentences else f"Point {i+1}"
        slide_content = ". ".join(sentaries[1:8])[:500]  # next 7 sentences
        
        if slide_content:
            slides.append({"title": slide_title, "content": slide_content})
    
    if not slides:
        slides = [{"title": title, "content": content[:1000]}]
    
    # Step 3: Generate PPT
    slides_result = native_slides(topic or title, slides)
    
    log("native_content_to_slides", "DONE", f"URL → {slides_result.get('slide_count', '?')} slides")
    
    return {
        "status": "ok",
        "source_url": url,
        "source_title": title,
        "slides_output": slides_result.get("output_path"),
        "slide_count": len(slides),
        "method": "FREE_NATIVE (web_fetch + python-pptx)"
    }


# =====================================================================
# 10. NATIVE TWITTER WRITER (Replace: felo-twitter-writer)
# =====================================================================

def native_twitter_writer(mode: str, account: str = None, topic: str = None, 
                          style_dna: str = None) -> dict:
    """
    Replace felo-twitter-writer with native approach.
    
    Mode 'analyze': Extract tweets from account → summarize style
    Mode 'write': Generate tweet from topic + style
    """
    log("native_twitter_writer", "ACTIVE", f"Mode: {mode}, Account: {account}")
    
    if mode == "analyze" and account:
        # Extract tweets from account
        tweets_result = native_x_search(f"from:{account.replace('@', '')}", limit=20)
        tweets = tweets_result.get("tweets", [])
        
        if not tweets:
            return {
                "status": "partial",
                "mode": "analyze",
                "message": "No tweets found for account",
                "method": "FREE_NATIVE"
            }
        
        # Summarize style (simple heuristic)
        total_len = sum(len(t.get("content", "")) for t in tweets)
        avg_len = total_len / len(tweets) if tweets else 0
        
        style_dna = f"""Style DNA: @{account}
- Average tweet length: {avg_len:.0f} chars
- Tweet count analyzed: {len(tweets)}
- Pattern: {'concise' if avg_len < 150 else 'detailed'}
"""
        
        return {
            "status": "ok",
            "mode": "analyze",
            "account": account,
            "style_dna": style_dna,
            "tweets_analyzed": len(tweets),
            "method": "FREE_NATIVE (native_x_search + heuristic)"
        }
    
    elif mode == "write" and topic:
        # Generate tweet (using MiniMax if available)
        prompt = f"Write a tweet (max 280 chars) about: {topic}"
        if style_dna:
            prompt += f"\nStyle: {style_dna}"
        
        # Try MiniMax
        try:
            minimax_key = os.environ.get("MINIMAX_API_KEY", "")
            if not minimax_key:
                cred_file = "/root/credential/api_key.json"
                if os.path.exists(cred_file):
                    with open(cred_file) as f:
                        creds = json.load(f)
                        if "minimax" in creds and creds["minimax"].get("keys"):
                            minimax_key = creds["minimax"]["keys"][0]
            
            if minimax_key:
                import requests
                payload = {
                    "model": "MiniMax-M2.7",
                    "messages": [{"role": "user", "content": prompt}],
                    "stream": False
                }
                resp = requests.post(
                    "https://api.minimax.io/v1/text/chatcompletion_v2",
                    headers={"Authorization": f"Bearer {minimax_key}"},
                    json=payload, timeout=30
                )
                tweet = resp.json()["choices"][0]["message"]["content"]
            else:
                tweet = f"[FREE] Tweet about {topic}: [MiniMax key needed for generation]"
        except Exception as e:
            tweet = f"[FREE] Tweet about {topic}: [Error: {e}]"
        
        return {
            "status": "ok",
            "mode": "write",
            "topic": topic,
            "tweet": tweet[:280],
            "method": "FREE_NATIVE (MiniMax-M2.7)"
        }
    
    return {"status": "error", "message": "Invalid mode or missing parameters"}


# =====================================================================
# MAIN / TEST
# =====================================================================

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="ILMA FELO FREE — Native Replacements")
    parser.add_argument("--skill", required=True, 
                        choices=["search", "x-search", "slides", "mindmap", 
                                "web-fetch", "youtube", "superagent", "kb-create", 
                                "kb-search", "content-slides", "twitter-writer"])
    parser.add_argument("--query")
    parser.add_argument("--url")
    parser.add_argument("--topic")
    parser.add_argument("--video-id")
    parser.add_argument("--video-url")
    parser.add_argument("--account")
    parser.add_argument("--mode")
    parser.add_argument("--kb-id")
    parser.add_argument("--style-dna")
    args = parser.parse_args()
    
    result = None
    
    if args.skill == "search":
        result = native_search(args.query or "AI trends")
    elif args.skill == "x-search":
        result = native_x_search(args.query or "AI")
    elif args.skill == "slides":
        result = native_slides(args.topic or "Demo Topic")
    elif args.skill == "mindmap":
        result = native_mindmap(args.topic or "Demo Mindmap")
    elif args.skill == "web-fetch":
        result = native_web_fetch(args.url or "https://example.com")
    elif args.skill == "youtube":
        result = native_youtube_subtitles(video_id=args.video_id, video_url=args.video_url)
    elif args.skill == "superagent":
        result = native_superagent(args.query or "Hello")
    elif args.skill == "kb-create":
        result = native_kb_create(args.topic or "New KB")
    elif args.skill == "kb-search":
        result = native_kb_search(args.kb_id or "default", args.query or "test")
    elif args.skill == "content-slides":
        result = native_content_to_slides(args.url or "https://example.com", args.topic)
    elif args.skill == "twitter-writer":
        result = native_twitter_writer(args.mode or "write", args.account, args.topic, args.style_dna)
    
    print(json.dumps(result, indent=2))
